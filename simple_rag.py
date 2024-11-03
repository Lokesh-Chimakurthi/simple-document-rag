import os
import csv
import numpy as np
from typing import List, Dict, Any, Tuple
import google.generativeai as genai
from qdrant_client import QdrantClient
from qdrant_client.http.models import Distance, VectorParams, PointStruct
import pathlib
import re
import PyPDF2
from docx import Document
from pptx import Presentation
import markdown

# Configure API key
genai.configure(api_key=os.getenv('GEMINI_API_KEY'))

class ChatHistory:
    def __init__(self):
        self.messages: List[Tuple[str, str]] = []  # List of (question, answer) tuples
    
    def add_interaction(self, question: str, answer: str):
        self.messages.append((question, answer))
    
    def get_history_text(self) -> str:
        return "\n".join([
            f"User: {msg[0]}\nAssistant: {msg[1]}"
            for msg in self.messages
        ])
    
    def clear(self):
        self.messages.clear()

class DocumentProcessor:
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        raise NotImplementedError

class CSVProcessor(DocumentProcessor):
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        documents = []
        with open(file_path, 'r') as file:
            reader = csv.DictReader(file)
            for row in reader:
                text = " ".join(f"{k}: {v}" for k, v in row.items())
                documents.append({
                    'content': text,
                    'metadata': {
                        'source': file_path,
                        'type': 'csv',
                        **row
                    }
                })
        return documents

class PDFProcessor(DocumentProcessor):
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        documents = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text.strip():
                    documents.append({
                        'content': text,
                        'metadata': {
                            'source': file_path,
                            'type': 'pdf',
                            'page': page_num + 1
                        }
                    })
        return documents

class WordProcessor(DocumentProcessor):
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        documents = []
        doc = Document(file_path)
        
        for para_num, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                documents.append({
                    'content': paragraph.text,
                    'metadata': {
                        'source': file_path,
                        'type': 'docx',
                        'paragraph': para_num + 1
                    }
                })
        return documents

class PowerPointProcessor(DocumentProcessor):
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        documents = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            if slide_text:
                documents.append({
                    'content': ' '.join(slide_text),
                    'metadata': {
                        'source': file_path,
                        'type': 'pptx',
                        'slide': slide_num + 1
                    }
                })
        return documents

class TextProcessor(DocumentProcessor):
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        documents = []
        with open(file_path, 'r', encoding='utf-8') as file:
            # Split text into paragraphs
            paragraphs = [p.strip() for p in file.read().split('\n\n')]
            
            for para_num, paragraph in enumerate(paragraphs):
                if paragraph:  # Skip empty paragraphs
                    documents.append({
                        'content': paragraph,
                        'metadata': {
                            'source': file_path,
                            'type': 'txt',
                            'paragraph': para_num + 1
                        }
                    })
        return documents

class MarkdownProcessor(DocumentProcessor):
    def __init__(self):
        self.md = markdown.Markdown(extensions=['extra'])
    
    def _clean_html(self, text: str) -> str:
        # Remove HTML tags while preserving content
        clean = re.compile('<.*?>')
        return re.sub(clean, '', text)
    
    def _split_sections(self, content: str) -> List[str]:
        # Split by headers and preserve structure
        sections = re.split(r'(#{1,6}\s.*\n)', content)
        result = []
        current_header = ""
        
        for section in sections:
            if section.strip():
                if re.match(r'#{1,6}\s', section):
                    current_header = section.strip()
                else:
                    text = section.strip()
                    if text:
                        if current_header:
                            result.append(f"{current_header}\n{text}")
                        else:
                            result.append(text)
        
        return result
    
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        documents = []
        
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
            sections = self._split_sections(content)
            
            for section_num, section in enumerate(sections):
                if section.strip():
                    # Convert markdown to HTML then clean it
                    html = self.md.convert(section)
                    cleaned_text = self._clean_html(html)
                    
                    if cleaned_text.strip():
                        documents.append({
                            'content': cleaned_text,
                            'metadata': {
                                'source': file_path,
                                'type': 'md',
                                'section': section_num + 1
                            }
                        })
        
        return documents

class DataExtractor:
    def __init__(self):
        self.processors = {
            '.csv': CSVProcessor(),
            '.pdf': PDFProcessor(),
            '.docx': WordProcessor(),
            '.doc': WordProcessor(),
            '.pptx': PowerPointProcessor(),
            '.ppt': PowerPointProcessor(),
            '.md': MarkdownProcessor(),
            '.markdown': MarkdownProcessor(),
            '.txt': TextProcessor(),
            '.text': TextProcessor()
        }

    def load_document(self, file_path: str) -> List[Dict[str, Any]]:
        ext = pathlib.Path(file_path).suffix.lower()
        processor = self.processors.get(ext)
        
        if not processor:
            raise ValueError(f"Unsupported file type: {ext}")
        
        return processor.process(file_path)

class VectorStoreManager:
    def __init__(self, db_path: str = 'qdrant_db'):
        self.client = QdrantClient(path=db_path)
        self.embedding_model = "models/text-embedding-004"

    def get_embedding(self, text: str) -> List[float]:
        try:
            response = genai.embed_content(
                model=self.embedding_model,
                content=text,
                task_type="retrieval_document"
            )
            if not response['embedding']:
                raise ValueError("Empty embedding received")
            return response['embedding']
        except Exception as e:
            print(f"Error getting embedding: {e}")
            return None

    def setup_collection(self, collection_name: str = "rag-demo") -> str:
        try:
            if self.client.collection_exists(collection_name):
                return collection_name
            
            self.client.create_collection(
                collection_name=collection_name,
                vectors_config=VectorParams(
                    size=768,
                    distance=Distance.COSINE
                )
            )
            return collection_name
        except Exception as e:
            print(f"Error setting up vector store: {e}")
            return None

    def index_documents(self, documents: List[Dict], collection_name: str):
        points = []
        for idx, doc in enumerate(documents):
            vector = self.get_embedding(doc['content'])
            if vector is None:
                continue
                
            point = PointStruct(
                id=idx,
                vector=vector,
                payload={
                    'content': doc['content'],
                    'metadata': doc['metadata']
                }
            )
            points.append(point)
        
        if points:
            self.client.upsert(
                collection_name=collection_name,
                points=points
            )

    def search_similar(self, query: str, collection_name: str, limit: int = 3):
        query_vector = self.get_embedding(query)
        if query_vector is None:
            return []
            
        try:
            results = self.client.search(
                collection_name=collection_name,
                query_vector=query_vector,
                limit=limit
            )
            return [hit.payload for hit in results]
        except Exception as e:
            print(f"Error during search: {e}")
            return []

class AnswerGenerator:
    def __init__(self):
        self.model_name = "gemini-1.5-flash-002"

    def generate_answer(self, query: str, context: str, chat_history: str) -> str:
        prompt = f"""
        Based on the following context and chat history, answer the question.
        If you cannot find the answer in the context, say "I don't have enough information."
        
        Previous Conversation:
        {chat_history}
        
        Context: {context}
        
        Current Question: {query}
        
        Answer:"""
        
        model = genai.GenerativeModel(self.model_name)
        response = model.generate_content(prompt)
        return response.text

class RAGOrchestrator:
    def __init__(self, db_path: str = 'qdrant_db'):
        self.extractor = DataExtractor()
        self.vector_store = VectorStoreManager(db_path)
        self.generator = AnswerGenerator()
        self.chat_history = ChatHistory()
        self.collection_name = None
        self.processed_files = []

    def setup_vector_store(self, collection_name: str = "rag-demo") -> str:
        self.collection_name = self.vector_store.setup_collection(collection_name)
        return self.collection_name

    def load_csv(self, file_path: str) -> List[Dict[str, Any]]:
        return self.extractor.load_document(file_path)

    def index_documents(self, documents: List[Dict], collection_name: str):
        self.vector_store.index_documents(documents, collection_name)

    def query(self, question: str, collection_name: str) -> Dict[str, Any]:
        similar_docs = self.vector_store.search_similar(question, collection_name)
        context = "\n".join([doc['content'] for doc in similar_docs])
        
        answer = self.generator.generate_answer(
            question, 
            context, 
            self.chat_history.get_history_text()
        )
        
        self.chat_history.add_interaction(question, answer)
        
        return {
            'question': question,
            'answer': answer,
            'sources': similar_docs
        }

    def clear_chat_history(self):
        self.chat_history.clear()

    def process_file_if_needed(self, file_path: str, collection_name: str) -> bool:
        if file_path in self.processed_files:
            print(f"File {file_path} already processed, skipping...")
            return True
            
        try:
            documents = self.extractor.load_document(file_path)
            if not documents:
                return False
                
            self.index_documents(documents, collection_name)
            self.processed_files.append(file_path)
            return True
        except Exception as e:
            print(f"Error processing file {file_path}: {e}")
            return False

def main():
    rag = RAGOrchestrator()
    
    collection_name = rag.setup_vector_store()
    if not collection_name:
        print("Failed to set up vector store")
        return
    
    # Example usage with different file types
    files = [
        'example.csv'
    ]
    
    for file_path in files:
        if not rag.process_file_if_needed(file_path, collection_name):
            print(f"Failed to process file: {file_path}")
            continue

    # Test queries with chat history
    questions = [
        "Which company does Lokesh work for?",
    ]
    
    for question in questions:
        result = rag.query(question, collection_name)
        print(f"\nQuestion: {result['question']}")
        print(f"Answer: {result['answer']}")
        print("\nSources:")
        for idx, source in enumerate(result['sources'], 1):
            print(f"\nSource {idx}:")
            print(source['content'][:200] + "...")

if __name__ == "__main__":
    main()